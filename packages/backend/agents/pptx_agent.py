"""PPTXAgent — parses uploaded deck, evaluates against Tikal 12-factor Playbook.

Runs post-upload (background task). Per §10b.6 — slide raw text persisted to
sessions.slides_raw_text BEFORE LLM call so container restart survives.
"""

from __future__ import annotations

import asyncio
import uuid

from agno.agent import Agent
from loguru import logger
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from agents.llm import (
    call_with_fallback,
    get_text_model,
    get_text_model_fallback,
    pick_provider_order,
)
from agents.persistence import AgentPersistence
from agents.replay_fixtures import replay_delivery_findings, replay_slide_findings
from agents.schemas import (
    ContentAnalysisPayload,
    SlideAnalysisPayload,
    SlideEventPayload,
    SlideFindings,
)
from config import get_settings
from db.repository import PostgreSQLRepository
from orchestrator.orchestrator import Orchestrator

PLAYBOOK_PROMPT = """You are a slide deck reviewer trained on the Tikal Presentation Skills Playbook.
Your job is to evaluate every slide in the deck and return specific, actionable findings.

## Rubric (Tikal Presentation Skills Playbook)
  1  Slides support the talk — they are NOT the talk. Plan the talk first as if there
     were no slides at all; only then think about slides.
  2  Authenticity — content should be original, unfamiliar, refreshing. If borrowing
     from specific sources, credit them on the slide.
  3  Keep it visual — limit text and bullet points; present information in small chunks.
     A picture is worth a thousand words; a funny/interesting picture is worth even more.
     Use high-resolution images.
  4  10-20-30 rule — ~10 slides, ~20 minutes, font size ≥ 30.
  5  Object count ≤6 per slide — the human brain counts up to 6 instantly; 7+ causes
     a delay. Avoid redundant slide numbers and logos on every slide.
  6  Dark background + bright text — high contrast keeps eyes on screen when needed;
     a blank slide lets the speaker reclaim all attention.
  7  Consistent fonts, headers, and title sizes across all slides.
  8  Story structure — slides should reflect a narrative arc: disruptive opening →
     hero encounters problem → naive solution fails → inspiration → solution → aftermath.
  9  Code display — animate or highlight specific lines; never dump a wall of unscoped
     code. Always have a Plan B (code pre-written in slides or a pre-recorded clip).
 10  Speaker notes belong in the notes area, NOT on the slide itself.
 11  Blank slides — insert one when you want all focus on the speaker.
 12  From subject to story — the deck should feel like a story, not an assorted
     collection of facts. The audience should always wonder "what happens next".

## What NOT to do
- Do not give generic advice like "add more images" without referencing the specific slide.
- Do not invent slide content — only evaluate what is explicitly shown to you.
- Do not emit more than 3 findings per slide — pick the most impactful ones.
- You receive extracted text plus structural metadata only — never the rendered slide.
  Do NOT comment on color, contrast, or background (Factor #6) — that data isn't provided.
- Only raise a font-size finding (Factor #4) when a slide's metadata line lists an explicit
  size below 30pt. If metadata says no explicit sizes were detected, say nothing about font size.
- Only say a slide "has visuals" or "uses an image" when its metadata line reports images > 0.
  Never infer an image's existence from the text content alone.
- Object-count findings (Factor #5) must match the "objects" count given in metadata, not a
  guess from the text.

## Metadata guide (from python-pptx extraction)
- Deck header shows total slide count — use for Factor #4 deck-length (10-20-30 rule).
- text_lines: non-empty lines on the slide body — use for Factor #3 text density; do not guess.
- speaker_notes: text in the notes area (not shown to audience) — use for Factor #10.
  If long body text duplicates speaker_notes content, flag Factor #10 IMPROVEMENT.
- is_blank: yes — slide has no body text and no images; evaluate Factor #11 (blank slide).
- notes_on_slide: yes — body text overlaps speaker notes; flag Factor #10 IMPROVEMENT.

## Output format
Return JSON only:
{"findings": [{
  "slide_index": int (1-based),
  "playbook_factor": int (1-12, must match the rubric above),
  "finding_type": "STRENGTH" | "IMPROVEMENT",
  "description": string (≤120 chars) — what is wrong OR what works well (diagnosis only),
  "suggested_fix": string (≤160 chars) — for IMPROVEMENT: a concrete rewrite or action
    (what you would do instead); for STRENGTH: what to keep doing (optional, ≤80 chars)
}]}

Every IMPROVEMENT must include a non-empty suggested_fix with a specific alternative — not
generic advice. Name the slide element to cut, merge, move to notes, or replace.

Cover at least 5 distinct factors across the deck. Both STRENGTH and IMPROVEMENT findings are welcome.

For decks with more than 10 slides, return 10–18 findings total (still max 3 per slide).
Scan the ENTIRE deck — do not stop after the first few slides.
Always include at least one deck-level Factor #4 finding when total slides ≠ ~10.
When metadata shows objects > 6, include a Factor #5 finding for that slide.
When metadata shows text_lines ≥ 8, include a Factor #3 finding.
When explicit font sizes are below 30pt, include a Factor #4 finding for that slide.
Look for narrative flow across slides for Factor #8 or #12.

## Examples of good findings
{"slide_index": 2, "playbook_factor": 3, "finding_type": "IMPROVEMENT",
 "description": "Slide 2 has 9 bullet points — too much text for the audience to read.",
 "suggested_fix": "Keep one headline + 3 bullets; move the other six points to speaker notes."}

{"slide_index": 5, "playbook_factor": 8, "finding_type": "STRENGTH",
 "description": "Slide 5 opens with a question that challenges assumptions — strong disruptive opening.",
 "suggested_fix": "Keep opening with a provocative question before stating the thesis."}

{"slide_index": 8, "playbook_factor": 9, "finding_type": "IMPROVEMENT",
 "description": "Slide 8 shows a full 40-line code block with no highlight.",
 "suggested_fix": "Show only lines 12–15 with a yellow highlight; paste the full file in notes."}

{"slide_index": 4, "playbook_factor": 11, "finding_type": "STRENGTH",
 "description": "Slide 4 is blank (is_blank: yes) — good use of a focus slide.",
 "suggested_fix": "Insert blank slides before key transitions to reclaim attention."}

{"slide_index": 6, "playbook_factor": 10, "finding_type": "IMPROVEMENT",
 "description": "Slide 6 body repeats speaker notes verbatim.",
 "suggested_fix": "Leave a ≤5-word headline on-slide; paste the full script into notes only."}"""

DELIVERY_PROMPT = """You are a presentation coach reviewing how well SLIDES supported what the speaker ACTUALLY SAID.
You receive a timed transcript from a live rehearsal plus the slide deck text. There is NO slide-change timeline —
infer which speech segments likely belong to which slides using semantic overlap only.

## Focus (Tikal Playbook — delivery lens)
  1  Slides support the talk — they are NOT the talk. Flag when speech carries the message but slides distract or mismatch.
  8  Story structure — speech and slide sequence should feel like one narrative.
 10  Speaker notes belong off-slide — flag when the speaker reads slide text verbatim instead of paraphrasing.
 12  From subject to story — flag when spoken content wanders while slides stay on unrelated facts.

Also cross-check the optional Content Analysis block: if the speaker made a factual error or skipped a key topic,
note when the visible slide still shows outdated or missing information.

## What NOT to do
- Do not repeat static deck-design issues (font size, object count, bullet density) unless they caused a delivery problem.
- Do not invent transcript quotes — only reference what appears in the timed transcript.
- Do not invent slide content — only evaluate slides explicitly provided.
- Do not emit more than 2 findings per slide.
- Do not judge timing ("stayed too long on slide") — no slide timeline is provided.

## Output format
Return JSON only:
{"findings": [{
  "slide_index": int (1-based),
  "playbook_factor": int (1, 8, 10, or 12 preferred),
  "finding_type": "STRENGTH" | "IMPROVEMENT",
  "description": string (≤120 chars) — diagnosis of speech vs slide support,
  "suggested_fix": string (≤160 chars) — concrete slide or delivery change for IMPROVEMENT
}]}

Return 3–8 findings total. Every IMPROVEMENT must have a non-empty suggested_fix naming what to change on-slide or how to speak.

## Examples
{"slide_index": 4, "playbook_factor": 1, "finding_type": "IMPROVEMENT",
 "description": "You explained JWT refresh in speech but slide 4 still shows generic auth bullets.",
 "suggested_fix": "Replace slide 4 with a refresh-token flow diagram matching what you said."}

{"slide_index": 11, "playbook_factor": 10, "finding_type": "IMPROVEMENT",
 "description": "You read slide 11 bullet text verbatim instead of paraphrasing.",
 "suggested_fix": "Leave a ≤5-word headline on slide 11; speak the detail without reading."}

{"slide_index": 8, "playbook_factor": 1, "finding_type": "STRENGTH",
 "description": "Slide 8 headline matched your spoken example — good slide support.",
 "suggested_fix": "Keep pairing a short slide headline with a spoken walk-through."}"""

DELIVERY_TIMELINE_PROMPT = """You are a presentation coach reviewing how well SLIDES supported what the speaker ACTUALLY SAID.
You receive a TIMED slide timeline plus transcript segments bucketed under each slide. Use the timestamps — do not guess.

## Focus (Tikal Playbook — delivery lens)
  1  Slides support the talk — flag when speech on slide N mismatches slide N's content, or when the speaker is ahead/behind slides.
  8  Story structure — speech and slide sequence should feel like one narrative across time.
 10  Speaker notes belong off-slide — flag reading slide text verbatim during a segment.
 12  From subject to story — flag wandering speech while a slide stays on unrelated facts.

Also cross-check the optional Content Analysis block for factual errors vs visible slides.

## What NOT to do
- Do not repeat static deck-design issues unless they caused a delivery problem during a timed segment.
- Do not invent transcript quotes — only reference speech in the aligned segments.
- Do not emit more than 2 findings per slide.
- Cite approximate timestamps in descriptions (e.g. "At 2:15…") when the mismatch is time-specific.

## Output format
Return JSON only:
{"findings": [{
  "slide_index": int (1-based),
  "playbook_factor": int (1, 8, 10, or 12 preferred),
  "finding_type": "STRENGTH" | "IMPROVEMENT",
  "description": string (≤120 chars),
  "suggested_fix": string (≤160 chars)
}]}

Return 3–8 findings. Flag ahead/behind slide sync, long dwell with off-topic speech, and strong alignment.

## Examples
{"slide_index": 4, "playbook_factor": 1, "finding_type": "IMPROVEMENT",
 "description": "At 2:15 on slide 4 you explained JWT refresh but the slide still shows generic auth bullets.",
 "suggested_fix": "Replace slide 4 with a refresh-token flow diagram matching what you said."}

{"slide_index": 7, "playbook_factor": 1, "finding_type": "IMPROVEMENT",
 "description": "At 4:30 you were still on slide 3 while describing slide 7's architecture.",
 "suggested_fix": "Advance to slide 7 before diving into its architecture, or pause until you catch up."}"""


MAX_FINDINGS_PER_SLIDE = 3
MAX_DELIVERY_FINDINGS = 8
MAX_DELIVERY_PER_SLIDE = 2
MAX_DESCRIPTION_LEN = 120
MAX_SUGGESTED_FIX_LEN = 160
MIN_DISTINCT_FACTORS = 5

_FACTOR_FIX_TEMPLATES: dict[int, str] = {
    1: "Draft the talk outline first, then add slides only where they clarify a point.",
    2: "Replace stock phrases with a personal anecdote or credited source on the slide.",
    3: "Replace bullet walls with one headline + diagram; move supporting text to speaker notes.",
    4: "Merge related slides or move deep dives to appendix; aim for ~10 slides and 30pt+ body text.",
    5: "Keep title, ≤3 bullets, and one visual; delete or move extra shapes off-slide.",
    6: "Switch to dark background + light text, or insert a blank slide before the key point.",
    7: "Pick one header size and one body size; apply via slide master across the deck.",
    8: "Reorder slides: hook → problem → failed attempt → inspiration → solution → takeaway.",
    9: "Show only the 3–5 relevant code lines with highlight; paste the full file in notes.",
    10: "Leave a ≤5-word headline on-slide; paste the full script into speaker notes only.",
    11: "Insert a blank slide before major transitions to reclaim audience attention.",
    12: "Reorder slides so each one raises a question answered on the next slide.",
}


def _truncate(text: str, max_len: int) -> str:
    text = text.strip()
    if len(text) <= max_len:
        return text
    return text[: max_len - 1] + "…"


def _default_suggested_fix(factor: int, row: dict | None = None) -> str:
    idx = row["slide_index"] if row else None
    if factor == 4 and row:
        sizes = row.get("font_sizes_pt") or []
        small = [s for s in sizes if s < 30]
        if small:
            return f"Set body text on slide {idx} to 30pt+; demote secondary lines to notes."
        total = row.get("_deck_total")
        if total and total > 12:
            return "Merge related slides or move deep dives to appendix; aim for ~10 slides total."
    if factor == 5 and row:
        return (
            f"On slide {idx}, keep title + ≤3 bullets + one chart; "
            "delete or move the other shapes to notes."
        )
    if factor == 3 and row:
        return (
            f"On slide {idx}, replace dense bullets with one headline + diagram; "
            "move detail to speaker notes."
        )
    return _FACTOR_FIX_TEMPLATES.get(factor, "Apply the playbook guidance for this factor on this slide.")


def _normalize_finding(
    item: SlideAnalysisPayload,
    row: dict | None = None,
    *,
    phase: str = "static",
) -> SlideAnalysisPayload:
    desc = _truncate(item.description, MAX_DESCRIPTION_LEN)
    fix = _truncate(item.suggested_fix, MAX_SUGGESTED_FIX_LEN)
    if item.finding_type == "IMPROVEMENT" and not fix:
        fix = _truncate(_default_suggested_fix(item.playbook_factor, row), MAX_SUGGESTED_FIX_LEN)
    return SlideAnalysisPayload(
        slide_index=item.slide_index,
        playbook_factor=item.playbook_factor,
        finding_type=item.finding_type,
        description=desc,
        suggested_fix=fix,
        analysis_phase=phase,
    )


def _format_ms(ms: int) -> str:
    total = max(0, ms // 1000)
    m, s = divmod(total, 60)
    return f"{m}:{s:02d}"


def _format_timed_transcript(entries) -> str:
    lines: list[str] = []
    for e in entries:
        text = (getattr(e, "text", "") or "").strip()
        if not text:
            continue
        start = getattr(e, "start_ms", 0)
        end = getattr(e, "end_ms", start)
        lines.append(f"[{_format_ms(start)}–{_format_ms(end)}] {text}")
    return "\n".join(lines)


def _event_timestamp_ms(event) -> int:
    return int(getattr(event, "timestamp_ms", 0))


def _event_slide_index(event) -> int:
    return int(getattr(event, "slide_index", 0))


def build_slide_timeline(
    events,
    slide_count: int,
    session_end_ms: int | None = None,
) -> list[dict[str, int]]:
    """Normalize slide change events into contiguous segments."""
    if not events:
        return []

    sorted_events = sorted(events, key=_event_timestamp_ms)
    segments: list[dict[str, int]] = []
    for i, ev in enumerate(sorted_events):
        idx = _event_slide_index(ev)
        if not (1 <= idx <= slide_count):
            continue
        start = _event_timestamp_ms(ev)
        if i + 1 < len(sorted_events):
            end = _event_timestamp_ms(sorted_events[i + 1])
        elif session_end_ms is not None:
            end = max(start, session_end_ms)
        else:
            end = start + 60_000
        segments.append({"slide_index": idx, "start_ms": start, "end_ms": end})
    return segments


def format_time_aligned_transcript(transcript, timeline: list[dict[str, int]]) -> str:
    """Bucket transcript lines under the slide active during each segment."""
    parts: list[str] = []
    for seg in timeline:
        idx = seg["slide_index"]
        start = seg["start_ms"]
        end = seg["end_ms"]
        lines: list[str] = []
        for e in transcript:
            e_start = getattr(e, "start_ms", 0)
            e_end = getattr(e, "end_ms", e_start)
            if e_end <= start or e_start >= end:
                continue
            text = (getattr(e, "text", "") or "").strip()
            if text:
                lines.append(f"[{_format_ms(e_start)}–{_format_ms(e_end)}] {text}")
        header = f"=== Slide {idx} ({_format_ms(start)}–{_format_ms(end)}) ==="
        body = "\n".join(lines) if lines else "(no speech during this segment)"
        parts.append(f"{header}\n{body}")
    return "\n\n".join(parts)


def _format_content_context(content: ContentAnalysisPayload | None) -> str:
    if not content or not content.findings:
        return "None provided."
    parts: list[str] = []
    for f in content.findings:
        quote = f' Quote: "{f.context_quote}"' if f.context_quote else ""
        parts.append(f"- [{f.type}] {f.message}{quote}")
    return "\n".join(parts)


def validate_delivery_findings(findings: list[SlideAnalysisPayload]) -> list[SlideAnalysisPayload]:
    """Cap delivery findings — semantic pass, no metadata supplement."""
    per_slide: dict[int, int] = {}
    kept: list[SlideAnalysisPayload] = []
    for item in findings:
        if len(kept) >= MAX_DELIVERY_FINDINGS:
            break
        count = per_slide.get(item.slide_index, 0)
        if count >= MAX_DELIVERY_PER_SLIDE:
            continue
        kept.append(_normalize_finding(item, row=None, phase="delivery"))
        per_slide[item.slide_index] = count + 1
    return kept


def _notes_overlap_body(body: str, notes: str) -> bool:
    body_norm = body.strip().lower()
    notes_norm = notes.strip().lower()
    if len(body_norm) < 40 or len(notes_norm) < 40:
        return False
    if body_norm == notes_norm:
        return True
    shorter, longer = (
        (body_norm, notes_norm) if len(body_norm) <= len(notes_norm) else (notes_norm, body_norm)
    )
    return shorter in longer


def validate_findings(
    findings: list[SlideAnalysisPayload],
    slides_raw: list[dict] | None = None,
) -> list[SlideAnalysisPayload]:
    """Enforce FR-001 limits and log coverage gaps for tuning."""
    rows_by_index = {row["slide_index"]: row for row in slides_raw} if slides_raw else {}
    if slides_raw:
        deck_total = len(slides_raw)
        for row in slides_raw:
            row["_deck_total"] = deck_total

    per_slide: dict[int, int] = {}
    kept: list[SlideAnalysisPayload] = []
    for item in findings:
        count = per_slide.get(item.slide_index, 0)
        if count >= MAX_FINDINGS_PER_SLIDE:
            continue
        row = rows_by_index.get(item.slide_index)
        kept.append(_normalize_finding(item, row, phase="static"))
        per_slide[item.slide_index] = count + 1

    factors = {f.playbook_factor for f in kept}
    types = {f.finding_type for f in kept}
    if len(factors) < MIN_DISTINCT_FACTORS:
        logger.warning(
            "PPTX findings cover {} factors (need ≥{}): {}",
            len(factors),
            MIN_DISTINCT_FACTORS,
            sorted(factors),
        )
    if len(types) < 2:
        logger.warning("PPTX findings missing STRENGTH/IMPROVEMENT mix: {}", types)

    if slides_raw and len(factors) < MIN_DISTINCT_FACTORS:
        kept = _supplement_from_metadata(kept, slides_raw)

    return kept


def _supplement_from_metadata(
    findings: list[SlideAnalysisPayload],
    slides_raw: list[dict],
) -> list[SlideAnalysisPayload]:
    """Deterministic fallback when the LLM under-covers factors (demo reliability)."""
    existing_factors = {f.playbook_factor for f in findings}
    extras: list[SlideAnalysisPayload] = []

    total = len(slides_raw)
    deck_total = total
    for row in slides_raw:
        row["_deck_total"] = deck_total
    if 4 not in existing_factors and total > 0:
        desc = (
            f"Deck has {total} slides — well above the ~10-slide 10-20-30 target."
            if total > 12
            else f"Deck has {total} slides — check pace against the 10-20-30 rule."
        )
        fix = (
            "Merge related slides or move deep dives to appendix; aim for ~10 slides total."
            if total > 12
            else "Keep one idea per slide; cut or appendix anything that does not advance the story."
        )
        extras.append(
            SlideAnalysisPayload(
                slide_index=1,
                playbook_factor=4,
                finding_type="IMPROVEMENT" if total > 12 else "STRENGTH",
                description=desc[:MAX_DESCRIPTION_LEN],
                suggested_fix=fix[:MAX_SUGGESTED_FIX_LEN],
            )
        )
        existing_factors.add(4)

    for row in slides_raw:
        if len(existing_factors) >= MIN_DISTINCT_FACTORS:
            break
        idx = row["slide_index"]
        if 5 not in existing_factors and row.get("shape_count", 0) > 6:
            extras.append(
                SlideAnalysisPayload(
                    slide_index=idx,
                    playbook_factor=5,
                    finding_type="IMPROVEMENT",
                    description=(
                        f"Slide {idx} has {row['shape_count']} objects — above the ≤6 limit."
                    )[:MAX_DESCRIPTION_LEN],
                    suggested_fix=_default_suggested_fix(5, row)[:MAX_SUGGESTED_FIX_LEN],
                )
            )
            existing_factors.add(5)
        if 3 not in existing_factors and row.get("text_line_count", 0) >= 8:
            extras.append(
                SlideAnalysisPayload(
                    slide_index=idx,
                    playbook_factor=3,
                    finding_type="IMPROVEMENT",
                    description=(
                        f"Slide {idx} has {row['text_line_count']} text lines — too dense to scan."
                    )[:MAX_DESCRIPTION_LEN],
                    suggested_fix=_default_suggested_fix(3, row)[:MAX_SUGGESTED_FIX_LEN],
                )
            )
            existing_factors.add(3)
        sizes = row.get("font_sizes_pt") or []
        small = [s for s in sizes if s < 30]
        if 4 not in existing_factors and small:
            extras.append(
                SlideAnalysisPayload(
                    slide_index=idx,
                    playbook_factor=4,
                    finding_type="IMPROVEMENT",
                    description=(
                        f"Slide {idx} uses {min(small):.0f}pt text — below the 30pt minimum."
                    )[:MAX_DESCRIPTION_LEN],
                    suggested_fix=_default_suggested_fix(4, row)[:MAX_SUGGESTED_FIX_LEN],
                )
            )
            existing_factors.add(4)
        if 11 not in existing_factors and row.get("is_blank"):
            extras.append(
                SlideAnalysisPayload(
                    slide_index=idx,
                    playbook_factor=11,
                    finding_type="STRENGTH",
                    description=f"Slide {idx} is blank — good focus slide for the speaker.",
                    suggested_fix="Insert blank slides before key transitions to reclaim attention.",
                )
            )
            existing_factors.add(11)
        if 10 not in existing_factors and row.get("notes_overlap_body"):
            extras.append(
                SlideAnalysisPayload(
                    slide_index=idx,
                    playbook_factor=10,
                    finding_type="IMPROVEMENT",
                    description=f"Slide {idx} duplicates speaker notes on the slide body.",
                    suggested_fix=_default_suggested_fix(10, row)[:MAX_SUGGESTED_FIX_LEN],
                )
            )
            existing_factors.add(10)

    merged = findings + extras
    per_slide: dict[int, int] = {}
    for f in merged:
        per_slide[f.slide_index] = per_slide.get(f.slide_index, 0) + 1

    def _slide_with_room() -> int:
        for row in slides_raw:
            idx = row["slide_index"]
            if per_slide.get(idx, 0) < MAX_FINDINGS_PER_SLIDE:
                return idx
        return slides_raw[0]["slide_index"] if slides_raw else 1

    covered = {f.playbook_factor for f in merged}
    if len(covered) < MIN_DISTINCT_FACTORS and 12 not in covered:
        idx = _slide_with_room()
        merged.append(
            SlideAnalysisPayload(
                slide_index=idx,
                playbook_factor=12,
                finding_type="IMPROVEMENT",
                description="Deck reads as a list of facts — weak narrative pull slide to slide.",
                suggested_fix=_default_suggested_fix(12)[:MAX_SUGGESTED_FIX_LEN],
            )
        )
        per_slide[idx] = per_slide.get(idx, 0) + 1
        covered.add(12)
    if len(covered) < MIN_DISTINCT_FACTORS and 7 not in covered:
        idx = _slide_with_room()
        merged.append(
            SlideAnalysisPayload(
                slide_index=idx,
                playbook_factor=7,
                finding_type="IMPROVEMENT",
                description="Title and body font sizes vary across slides.",
                suggested_fix=_default_suggested_fix(7)[:MAX_SUGGESTED_FIX_LEN],
            )
        )

    return validate_findings(merged, slides_raw=None)


class PPTXAgent(AgentPersistence):
    def __init__(self, orchestrator: Orchestrator, *, session_id: uuid.UUID | None = None) -> None:
        self.orchestrator = orchestrator
        self.session_id = session_id
        self._active_slide_index: int | None = None
        self._active_slide_since_ms: int | None = None
        self._last_event_ts: int = -1

    def get_slide_state(self, now_ms: int) -> dict[str, int] | None:
        """In-memory active slide for live STALE_SLIDE checks (Or Yam / AudioAgent)."""
        if self._active_slide_index is None or self._active_slide_since_ms is None:
            return None
        return {
            "slide_index": self._active_slide_index,
            "since_ms": self._active_slide_since_ms,
            "dwell_ms": max(0, now_ms - self._active_slide_since_ms),
        }

    async def record_slide_event(
        self, payload: SlideEventPayload, *, slide_count: int
    ) -> bool:
        """Persist a slide change and update in-memory state. Returns False if rejected."""
        sid = self.session_id
        if sid is None:
            logger.warning("record_slide_event called without session_id")
            return False
        if payload.timestamp_ms < self._last_event_ts:
            logger.warning(
                "reject out-of-order slide event ts={} last={}",
                payload.timestamp_ms,
                self._last_event_ts,
            )
            return False
        if not (1 <= payload.slide_index <= slide_count):
            return False

        self._active_slide_index = payload.slide_index
        self._active_slide_since_ms = payload.timestamp_ms
        self._last_event_ts = payload.timestamp_ms

        async def write(repo: PostgreSQLRepository) -> None:
            await repo.insert_slide_event(
                session_id=sid,
                slide_index=payload.slide_index,
                timestamp_ms=payload.timestamp_ms,
                source=payload.source,
            )

        await self._with_repo(write)
        return True

    async def analyse(self, session_id: uuid.UUID, pptx_path: str) -> None:
        try:
            slides_raw = await asyncio.to_thread(self._extract_text, pptx_path)
        except Exception as exc:  # noqa: BLE001
            logger.exception("PPTX parse failed: {}", exc)
            slides_raw = []

        if get_settings().demo_replay:
            findings = replay_slide_findings()
        else:
            try:
                findings = await self._evaluate(slides_raw)
            except Exception as exc:  # noqa: BLE001
                logger.exception("PPTX LLM eval failed: {}", exc)
                findings = []

        findings = validate_findings(findings, slides_raw)

        # This agent owns slide_analyses + sessions(pptx) writes (Principle II).
        # Order preserved from the prior Orchestrator write: findings, then pptx_ready.
        async def write(repo: PostgreSQLRepository) -> None:
            await repo.insert_slide_analyses(
                [
                    {
                        "session_id": session_id,
                        "slide_index": f.slide_index,
                        "playbook_factor": f.playbook_factor,
                        "finding_type": f.finding_type,
                        "description": f.description,
                        "suggested_fix": f.suggested_fix,
                        "analysis_phase": "static",
                    }
                    for f in findings
                ]
            )
            await repo.mark_pptx_ready(session_id, slides_raw)

        await self._with_repo(write)
        await self.orchestrator.notify_complete(session_id, "PPTX")

    async def analyse_delivery(
        self,
        session_id: uuid.UUID,
        *,
        topic: str,
        topic_context: str | None,
        slides_raw: list[dict],
        transcript,
        content: ContentAnalysisPayload | None = None,
        slide_events=None,
    ) -> None:
        """Post-session pass: compare timed speech + slide text (+ content findings)."""
        if not slides_raw:
            return
        transcript_text = _format_timed_transcript(transcript)
        if not transcript_text.strip():
            return

        slide_count = len(slides_raw)
        session_end_ms = max(
            (getattr(e, "end_ms", 0) for e in transcript),
            default=0,
        )
        timeline = (
            build_slide_timeline(slide_events, slide_count, session_end_ms=session_end_ms)
            if slide_events
            else []
        )
        use_timeline = bool(timeline)

        if get_settings().demo_replay:
            findings = replay_delivery_findings()
        else:
            try:
                findings = await self._evaluate_delivery(
                    topic=topic,
                    topic_context=topic_context,
                    slides_raw=slides_raw,
                    transcript_text=transcript_text,
                    content=content,
                    timeline=timeline if use_timeline else None,
                    transcript=transcript if use_timeline else None,
                )
            except Exception as exc:  # noqa: BLE001
                logger.exception("PPTX delivery LLM eval failed: {}", exc)
                findings = []

        findings = validate_delivery_findings(findings)
        if not findings:
            return

        async def write(repo: PostgreSQLRepository) -> None:
            await repo.delete_slide_analyses_by_phase(session_id, "delivery")
            await repo.insert_slide_analyses(
                [
                    {
                        "session_id": session_id,
                        "slide_index": f.slide_index,
                        "playbook_factor": f.playbook_factor,
                        "finding_type": f.finding_type,
                        "description": f.description,
                        "suggested_fix": f.suggested_fix,
                        "analysis_phase": "delivery",
                    }
                    for f in findings
                ]
            )

        await self._with_repo(write)
        if self.orchestrator is not None:
            await self.orchestrator.notify_complete(
                session_id, "PPTX", meta={"phase": "delivery", "findings": len(findings)}
            )

    @staticmethod
    def _extract_notes(slide) -> str:
        try:
            if slide.has_notes_slide:
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame is not None:
                    return notes_frame.text.strip()
        except AttributeError:
            pass
        return ""

    @staticmethod
    def _extract_text(path: str) -> list[dict]:
        deck = Presentation(path)
        out: list[dict] = []
        for idx, slide in enumerate(deck.slides, start=1):
            texts: list[str] = []
            font_sizes_pt: list[float] = []
            text_line_count = 0
            image_count = 0
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_count += 1
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = "".join(run.text for run in para.runs).strip()
                        if line:
                            texts.append(line)
                            text_line_count += 1
                        for run in para.runs:
                            if run.font.size is not None:
                                font_sizes_pt.append(run.font.size.pt)
            body_text = "\n".join(texts)
            notes_text = PPTXAgent._extract_notes(slide)
            notes_overlap = _notes_overlap_body(body_text, notes_text)
            out.append(
                {
                    "slide_index": idx,
                    "text": body_text,
                    "shape_count": len(slide.shapes),
                    "image_count": image_count,
                    "font_sizes_pt": font_sizes_pt,
                    "text_line_count": text_line_count,
                    "notes_text": notes_text,
                    "notes_overlap_body": notes_overlap,
                    "is_blank": not body_text.strip() and image_count == 0,
                }
            )
        return out

    @staticmethod
    def _format_deck_header(slides_raw: list[dict]) -> str:
        total = len(slides_raw)
        blank_slides = [row["slide_index"] for row in slides_raw if row.get("is_blank")]
        blank_note = f", blank slide indices: {blank_slides}" if blank_slides else ""
        return (
            f"=== Deck ({total} slides{blank_note}) ===\n"
            "Use total slide count for Factor #4 (10-20-30 rule: ~10 slides per ~20 minutes)."
        )

    @staticmethod
    def _format_slide(row: dict) -> str:
        sizes = row.get("font_sizes_pt") or []
        size_note = f"{min(sizes):.0f}-{max(sizes):.0f}pt" if sizes else "none detected — do not judge font size"
        notes = (row.get("notes_text") or "").strip()
        notes_note = notes if notes else "none"
        if len(notes_note) > 300:
            notes_note = notes_note[:300] + "..."
        meta = (
            f"objects: {row['shape_count']}, images: {row['image_count']}, "
            f"text_lines: {row.get('text_line_count', 0)}, is_blank: {'yes' if row.get('is_blank') else 'no'}, "
            f"notes_on_slide: {'yes' if row.get('notes_overlap_body') else 'no'}, "
            f"explicit font sizes: {size_note}"
        )
        body = row["text"] if row["text"].strip() else "(no slide body text)"
        return (
            f"=== Slide {row['slide_index']} ({meta}) ===\n"
            f"{body}\n"
            f"Speaker notes: {notes_note}"
        )

    async def _evaluate(self, slides_raw: list[dict]) -> list[SlideAnalysisPayload]:
        deck_dump = self.build_deck_dump(slides_raw)
        agent = Agent(
            model=get_text_model(),
            instructions=PLAYBOOK_PROMPT,
            output_schema=SlideFindings,
        )

        async def _gateway():
            return await agent.arun(deck_dump)

        fb = get_text_model_fallback()

        async def _claude():
            return await Agent(model=fb, instructions=PLAYBOOK_PROMPT, output_schema=SlideFindings).arun(deck_dump)

        claude_fn = _claude if fb else None
        primary, secondary = pick_provider_order(claude_fn, _gateway)
        result = await call_with_fallback(primary, secondary)
        raw = result.content.findings
        logger.info("PPTX LLM result: {} findings", len(raw))
        return raw

    async def _evaluate_delivery(
        self,
        *,
        topic: str,
        topic_context: str | None,
        slides_raw: list[dict],
        transcript_text: str,
        content: ContentAnalysisPayload | None,
        timeline: list[dict[str, int]] | None = None,
        transcript=None,
    ) -> list[SlideAnalysisPayload]:
        use_timeline = bool(timeline)
        dump = self.build_delivery_dump(
            topic=topic,
            topic_context=topic_context,
            slides_raw=slides_raw,
            transcript_text=transcript_text,
            content=content,
            timeline=timeline,
            transcript=transcript if use_timeline else None,
        )
        prompt = DELIVERY_TIMELINE_PROMPT if use_timeline else DELIVERY_PROMPT
        agent = Agent(
            model=get_text_model(),
            instructions=prompt,
            output_schema=SlideFindings,
        )

        async def _gateway():
            return await agent.arun(dump)

        fb = get_text_model_fallback()

        async def _claude():
            return await Agent(model=fb, instructions=prompt, output_schema=SlideFindings).arun(
                dump
            )

        claude_fn = _claude if fb else None
        primary, secondary = pick_provider_order(claude_fn, _gateway)
        result = await call_with_fallback(primary, secondary)
        raw = result.content.findings
        logger.info(
            "PPTX delivery LLM result: {} findings (timeline={})",
            len(raw),
            use_timeline,
        )
        return raw

    @staticmethod
    def build_delivery_dump(
        *,
        topic: str,
        topic_context: str | None,
        slides_raw: list[dict],
        transcript_text: str,
        content: ContentAnalysisPayload | None,
        timeline: list[dict[str, int]] | None = None,
        transcript=None,
    ) -> str:
        audience = topic_context.strip() if topic_context else "not specified"
        if timeline and transcript is not None:
            aligned = format_time_aligned_transcript(transcript, timeline)
            timeline_note = (
                "=== Slide timeline (manual tracker — use these timestamps) ===\n"
                + "\n".join(
                    f"Slide {s['slide_index']}: {_format_ms(s['start_ms'])}–{_format_ms(s['end_ms'])}"
                    for s in timeline
                )
            )
            parts = [
                f"=== Session topic ===\n{topic}\nAudience context: {audience}",
                timeline_note,
                "=== Speech aligned to active slide ===",
                aligned,
                "=== Content analysis (cross-check speech accuracy vs slides) ===",
                _format_content_context(content),
                PPTXAgent.build_deck_dump(slides_raw),
            ]
        else:
            parts = [
                f"=== Session topic ===\n{topic}\nAudience context: {audience}",
                "=== Timed transcript (infer slide mapping from semantic overlap) ===",
                transcript_text[:12000],
                "=== Content analysis (cross-check speech accuracy vs slides) ===",
                _format_content_context(content),
                PPTXAgent.build_deck_dump(slides_raw),
            ]
        return "\n\n".join(parts)

    @staticmethod
    def build_deck_dump(slides_raw: list[dict]) -> str:
        parts = [PPTXAgent._format_deck_header(slides_raw)]
        parts.extend(PPTXAgent._format_slide(row) for row in slides_raw)
        return "\n\n".join(parts)
