"""Tests for PPTXAgent — extraction, validation, and deck dump (plan steps 1, 2, 6)."""

from __future__ import annotations

import tempfile
import uuid
from pathlib import Path
from unittest.mock import AsyncMock, MagicMock, patch

import pytest
from pptx import Presentation

from agents.pptx_agent import (
    MAX_DELIVERY_FINDINGS,
    MIN_DISTINCT_FACTORS,
    PPTXAgent,
    SlideAnalysisPayload,
    build_slide_timeline,
    format_time_aligned_transcript,
    validate_delivery_findings,
    validate_findings,
    _format_timed_transcript,
    _notes_overlap_body,
    _supplement_from_metadata,
)
from agents.replay_fixtures import replay_delivery_findings, replay_slide_events
from agents.schemas import SlideEventPayload

REPO_ROOT = Path(__file__).resolve().parents[3]
REAT_DECK = REPO_ROOT / "docs" / "testing" / "reat-19.2-news-example.pptx"


def _write_deck(path: Path) -> None:
    prs = Presentation()
    slide1 = prs.slides.add_slide(prs.slide_layouts[1])
    slide1.shapes.title.text = "Title"
    slide1.placeholders[1].text = "Line one\nLine two\nLine three"
    slide1.notes_slide.notes_text_frame.text = "Speaker note for slide 1"
    prs.slides.add_slide(prs.slide_layouts[6])
    prs.save(str(path))


def test_notes_overlap_body():
    body = "This is a long paragraph that appears on the slide body for the audience to read."
    assert _notes_overlap_body(body, body) is True
    assert _notes_overlap_body("short", "other notes") is False


def test_extract_text_metadata():
    with tempfile.TemporaryDirectory() as tmp:
        pptx_path = Path(tmp) / "test.pptx"
        _write_deck(pptx_path)
        rows = PPTXAgent._extract_text(str(pptx_path))

    assert len(rows) == 2
    assert rows[0]["text_line_count"] >= 2
    assert rows[0]["notes_text"] == "Speaker note for slide 1"
    assert rows[0]["is_blank"] is False
    assert rows[1]["is_blank"] is True


@pytest.mark.skipif(not REAT_DECK.exists(), reason="reat test deck missing")
def test_extract_reat_example_deck():
    rows = PPTXAgent._extract_text(str(REAT_DECK))
    assert len(rows) == 22
    assert all("text_line_count" in r for r in rows)
    assert all("notes_text" in r for r in rows)
    assert all("notes_overlap_body" in r for r in rows)
    heavy = max(rows, key=lambda r: r["text_line_count"])
    assert heavy["text_line_count"] >= 10


def test_format_deck_header_includes_slide_count():
    slides_raw = [{"slide_index": 1, "is_blank": False}, {"slide_index": 2, "is_blank": True}]
    header = PPTXAgent._format_deck_header(slides_raw)
    assert "2 slides" in header
    assert "Factor #4" in header


def test_format_slide_includes_notes_and_text_lines():
    row = {
        "slide_index": 3,
        "text": "Bullet A\nBullet B",
        "shape_count": 4,
        "image_count": 1,
        "text_line_count": 2,
        "font_sizes_pt": [24.0],
        "notes_text": "Remember to pause here",
        "notes_overlap_body": False,
        "is_blank": False,
    }
    formatted = PPTXAgent._format_slide(row)
    assert "text_lines: 2" in formatted
    assert "Speaker notes: Remember to pause here" in formatted
    assert "notes_on_slide: no" in formatted


def test_validate_findings_caps_per_slide_and_description():
    findings = [
        SlideAnalysisPayload(slide_index=1, playbook_factor=3, finding_type="IMPROVEMENT", description="a" * 200),
        SlideAnalysisPayload(slide_index=1, playbook_factor=4, finding_type="IMPROVEMENT", description="b"),
        SlideAnalysisPayload(slide_index=1, playbook_factor=5, finding_type="IMPROVEMENT", description="c"),
        SlideAnalysisPayload(slide_index=1, playbook_factor=7, finding_type="IMPROVEMENT", description="d"),
    ]
    kept = validate_findings(findings)
    assert len(kept) == 3
    assert len(kept[0].description) <= 120


def test_improvement_without_fix_gets_backfill():
    findings = [
        SlideAnalysisPayload(
            slide_index=2,
            playbook_factor=5,
            finding_type="IMPROVEMENT",
            description="Too many shapes on slide 2.",
        )
    ]
    kept = validate_findings(findings)
    assert kept[0].suggested_fix
    assert "Instead" not in kept[0].description


def test_supplement_from_metadata_reaches_five_factors():
    slides_raw = [
        {
            "slide_index": 1,
            "shape_count": 9,
            "text_line_count": 12,
            "font_sizes_pt": [18.0],
            "is_blank": False,
            "notes_overlap_body": False,
        },
        {
            "slide_index": 2,
            "shape_count": 2,
            "text_line_count": 1,
            "font_sizes_pt": [],
            "is_blank": False,
            "notes_overlap_body": False,
        },
    ]
    base = [
        SlideAnalysisPayload(
            slide_index=2,
            playbook_factor=8,
            finding_type="STRENGTH",
            description="Strong story opening.",
        )
    ]
    out = _supplement_from_metadata(base, slides_raw)
    factors = {f.playbook_factor for f in out}
    assert len(factors) >= MIN_DISTINCT_FACTORS


@pytest.mark.skipif(not REAT_DECK.exists(), reason="reat test deck missing")
def test_build_deck_dump_reat():
    rows = PPTXAgent._extract_text(str(REAT_DECK))
    dump = PPTXAgent.build_deck_dump(rows)
    assert "=== Deck (22 slides" in dump
    assert "Speaker notes:" in dump


@pytest.mark.asyncio
async def test_analyse_persists_and_notifies():
    orch = AsyncMock()
    agent = PPTXAgent(orch)
    session_id = uuid.UUID("00000000-0000-0000-0000-000000000001")
    slides_raw = [
        {
            "slide_index": 1,
            "text": "Hi",
            "shape_count": 2,
            "image_count": 0,
            "font_sizes_pt": [],
            "text_line_count": 1,
            "notes_text": "",
            "notes_overlap_body": False,
            "is_blank": False,
        }
    ]
    mock_findings = [
        SlideAnalysisPayload(
            slide_index=1,
            playbook_factor=3,
            finding_type="IMPROVEMENT",
            description="Slide 1 needs a visual anchor.",
            suggested_fix="Replace text with one headline + diagram.",
        )
    ]
    written: list[dict] = []

    async def capture_write(fn):
        repo = AsyncMock()
        repo.insert_slide_analyses = AsyncMock(side_effect=lambda items: written.extend(items))
        repo.mark_pptx_ready = AsyncMock()
        await fn(repo)
        return None

    with patch.object(agent, "_extract_text", return_value=slides_raw), patch.object(
        agent, "_evaluate", new=AsyncMock(return_value=mock_findings)
    ), patch.object(agent, "_with_repo", side_effect=capture_write), patch(
        "agents.pptx_agent.get_settings"
    ) as mock_settings:
        mock_settings.return_value.demo_replay = False
        await agent.analyse(session_id, "/tmp/x.pptx")

    assert len(written) >= 1
    assert written[0]["suggested_fix"]
    assert written[0]["analysis_phase"] == "static"
    orch.notify_complete.assert_awaited_once_with(session_id, "PPTX")


def test_format_timed_transcript():
    class _Entry:
        def __init__(self, start_ms, end_ms, text):
            self.start_ms = start_ms
            self.end_ms = end_ms
            self.text = text

    out = _format_timed_transcript([_Entry(0, 2500, "Hello world"), _Entry(65000, 67000, "Done")])
    assert "[0:00–0:02] Hello world" in out
    assert "[1:05–1:07] Done" in out


def test_validate_delivery_findings_caps():
    findings = [
        SlideAnalysisPayload(
            slide_index=i + 1,
            playbook_factor=1,
            finding_type="IMPROVEMENT",
            description=f"issue {i}",
            suggested_fix=f"fix {i}",
        )
        for i in range(10)
    ]
    kept = validate_delivery_findings(findings)
    assert len(kept) == MAX_DELIVERY_FINDINGS
    assert all(f.analysis_phase == "delivery" for f in kept)


def test_replay_delivery_findings_have_phase():
    findings = replay_delivery_findings()
    assert len(findings) >= 2
    assert all(f.analysis_phase == "delivery" for f in findings)
    for f in findings:
        if f.finding_type == "IMPROVEMENT":
            assert f.suggested_fix.strip()


@pytest.mark.asyncio
async def test_analyse_delivery_skips_without_transcript():
    orch = AsyncMock()
    agent = PPTXAgent(orch)
    session_id = uuid.uuid4()

    class _Session:
        topic = "React 19"
        topic_context = None
        slides_raw_text = [{"slide_index": 1, "text": "Hi"}]

    mock_repo = AsyncMock()
    mock_repo.get_session = AsyncMock(return_value=_Session())
    mock_repo.read_transcript = AsyncMock(return_value=[])

    class _Ctx:
        async def __aenter__(self):
            return object()

        async def __aexit__(self, *args):
            return None

    with patch("agents.pptx_agent.get_session_maker") as mock_maker, patch(
        "agents.pptx_agent.PostgreSQLRepository", return_value=mock_repo
    ):
        mock_maker.return_value = MagicMock(return_value=_Ctx())
        with patch.object(agent, "_with_repo", new=AsyncMock()) as mock_repo_write:
            await agent.analyse_delivery(session_id)
    mock_repo_write.assert_not_awaited()


@pytest.mark.asyncio
async def test_analyse_delivery_persists_delivery_phase():
    orch = AsyncMock()
    agent = PPTXAgent(orch)
    session_id = uuid.uuid4()
    deleted: list[str] = []
    written: list[dict] = []

    async def capture_write(fn):
        repo = AsyncMock()
        repo.delete_slide_analyses_by_phase = AsyncMock(
            side_effect=lambda _sid, phase: deleted.append(phase)
        )
        repo.insert_slide_analyses = AsyncMock(side_effect=lambda items: written.extend(items))
        await fn(repo)

    class _Entry:
        start_ms = 0
        end_ms = 2000
        text = "Today we cover authentication flows."

    class _Session:
        topic = "Auth in React"
        topic_context = "Senior devs"
        slides_raw_text = [{"slide_index": 1, "text": "Auth overview"}]

    mock_repo = AsyncMock()
    mock_repo.get_session = AsyncMock(return_value=_Session())
    mock_repo.read_transcript = AsyncMock(return_value=[_Entry()])

    class _Ctx:
        async def __aenter__(self):
            return object()

        async def __aexit__(self, *args):
            return None

    with patch("agents.pptx_agent.get_session_maker") as mock_maker, patch(
        "agents.pptx_agent.PostgreSQLRepository", return_value=mock_repo
    ), patch.object(agent, "_with_repo", side_effect=capture_write), patch(
        "agents.pptx_agent.get_settings"
    ) as mock_settings:
        mock_maker.return_value = MagicMock(return_value=_Ctx())
        mock_settings.return_value.demo_replay = True
        await agent.analyse_delivery(session_id)

    assert deleted == ["delivery"]
    assert written
    assert all(w["analysis_phase"] == "delivery" for w in written)
    orch.notify_complete.assert_awaited_once()
    assert orch.notify_complete.await_args.args[1] == "PPTX"


class _StubEvent:
    def __init__(self, slide_index: int, timestamp_ms: int):
        self.slide_index = slide_index
        self.timestamp_ms = timestamp_ms


class _StubTranscript:
    def __init__(self, start_ms: int, end_ms: int, text: str):
        self.start_ms = start_ms
        self.end_ms = end_ms
        self.text = text


def test_build_slide_timeline_segments():
    events = [_StubEvent(1, 0), _StubEvent(3, 45000), _StubEvent(2, 90000)]
    timeline = build_slide_timeline(events, slide_count=5, session_end_ms=120000)
    assert len(timeline) == 3
    assert timeline[0] == {"slide_index": 1, "start_ms": 0, "end_ms": 45000}
    assert timeline[1] == {"slide_index": 3, "start_ms": 45000, "end_ms": 90000}
    assert timeline[2] == {"slide_index": 2, "start_ms": 90000, "end_ms": 120000}


def test_format_time_aligned_transcript_buckets_speech():
    timeline = [{"slide_index": 1, "start_ms": 0, "end_ms": 5000}]
    transcript = [
        _StubTranscript(0, 2000, "Opening"),
        _StubTranscript(6000, 8000, "Later"),
    ]
    aligned = format_time_aligned_transcript(transcript, timeline)
    assert "Opening" in aligned
    assert "Later" not in aligned
    assert "Slide 1" in aligned


def test_build_delivery_dump_uses_timeline_when_provided():
    timeline = [{"slide_index": 1, "start_ms": 0, "end_ms": 5000}]
    transcript = [_StubTranscript(0, 2000, "Hello")]
    dump = PPTXAgent.build_delivery_dump(
        topic="Test",
        topic_context=None,
        slides_raw=[
            {
                "slide_index": 1,
                "text": "Title",
                "shape_count": 2,
                "image_count": 0,
                "text_line_count": 1,
                "notes_text": "",
                "notes_overlap_body": False,
                "is_blank": False,
            }
        ],
        transcript_text="unused",
        content=None,
        timeline=timeline,
        transcript=transcript,
    )
    assert "Slide timeline" in dump
    assert "Speech aligned to active slide" in dump
    assert "infer slide mapping" not in dump


@pytest.mark.asyncio
async def test_record_slide_event_updates_state_and_persists():
    orch = AsyncMock()
    agent = PPTXAgent(orch, session_id=uuid.uuid4())
    written: list[dict] = []

    async def capture_write(fn):
        repo = AsyncMock()
        repo.insert_slide_event = AsyncMock(side_effect=lambda **kw: written.append(kw))
        await fn(repo)

    with patch.object(agent, "_with_repo", side_effect=capture_write):
        ok = await agent.record_slide_event(
            SlideEventPayload(slide_index=2, timestamp_ms=1000), slide_count=10
        )
        assert ok is True
        rejected = await agent.record_slide_event(
            SlideEventPayload(slide_index=1, timestamp_ms=500), slide_count=10
        )
        assert rejected is False

    assert written[0]["slide_index"] == 2
    state = agent.get_slide_state(5000)
    assert state == {"slide_index": 2, "since_ms": 1000, "dwell_ms": 4000}


def test_replay_slide_events_fixture():
    events = replay_slide_events()
    assert len(events) >= 2
    assert _StubEvent(1, 0).slide_index == events[0].slide_index


@pytest.mark.asyncio
async def test_analyse_delivery_with_slide_events_uses_replay():
    orch = AsyncMock()
    session_id = uuid.uuid4()
    agent = PPTXAgent(orch, session_id=session_id)
    written: list[dict] = []

    async def capture_write(fn):
        repo = AsyncMock()
        repo.delete_slide_analyses_by_phase = AsyncMock()
        repo.insert_slide_analyses = AsyncMock(side_effect=lambda items: written.extend(items))
        await fn(repo)

    class _Session:
        topic = "Auth"
        topic_context = None
        slides_raw_text = [{"slide_index": i, "text": f"S{i}"} for i in range(1, 6)]

    mock_repo = AsyncMock()
    mock_repo.get_session = AsyncMock(return_value=_Session())
    mock_repo.read_transcript = AsyncMock(return_value=[_StubTranscript(0, 2000, "Intro")])

    class _Ctx:
        async def __aenter__(self):
            return object()

        async def __aexit__(self, *args):
            return None

    events = replay_slide_events()
    with patch("agents.pptx_agent.get_session_maker") as mock_maker, patch(
        "agents.pptx_agent.PostgreSQLRepository", return_value=mock_repo
    ), patch.object(agent, "_with_repo", side_effect=capture_write), patch(
        "agents.pptx_agent.get_settings"
    ) as mock_settings:
        mock_maker.return_value = MagicMock(return_value=_Ctx())
        mock_settings.return_value.demo_replay = True
        await agent.analyse_delivery(session_id, slide_events=events)

    assert written
    assert all(w["analysis_phase"] == "delivery" for w in written)


@pytest.mark.live
@pytest.mark.skipif(not REAT_DECK.exists(), reason="reat test deck missing")
@pytest.mark.asyncio
async def test_live_eval_reat_deck():
    """Run with: pytest -m live tests/test_pptx_agent.py -v (requires LITELLM_API_KEY, DEMO_MODE unset)."""
    rows = PPTXAgent._extract_text(str(REAT_DECK))
    agent = PPTXAgent(AsyncMock())
    findings = await agent._evaluate(rows)
    validated = validate_findings(findings, rows)
    factors = {f.playbook_factor for f in validated}
    types = {f.finding_type for f in validated}
    assert len(validated) >= 5
    assert len(factors) >= MIN_DISTINCT_FACTORS
    assert types >= {"STRENGTH", "IMPROVEMENT"}
    for f in validated:
        if f.finding_type == "IMPROVEMENT":
            assert f.suggested_fix.strip(), f"factor {f.playbook_factor} missing fix"
